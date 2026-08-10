"""
PPT 生成工具 - 基于 python-pptx 生成 PowerPoint 演示文稿。
"""
import base64
import json
import logging
import os
import uuid
from io import BytesIO
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from app.core.tools import BaseTool, BaseToolParam, ToolRegistry

logger = logging.getLogger(__name__)

# PPT 文件缓存 key 前缀
PPT_CACHE_PREFIX = "ppt_file:"
# 缓存过期时间（秒），默认 24 小时
PPT_CACHE_EXPIRE = 86400


def save_ppt_to_cache(file_id: str, file_base64: str, file_name: str) -> bool:
    """将 PPT 文件存入 Redis 缓存"""
    try:
        from app.database.redis_utils import redis_utils
        return redis_utils.set(
            f"{PPT_CACHE_PREFIX}{file_id}",
            json.dumps({"base64": file_base64, "file_name": file_name}, ensure_ascii=False),
            exp=PPT_CACHE_EXPIRE
        )
    except Exception as e:
        logger.error(f"PPT 存入 Redis 失败: {e}", exc_info=True)
        return False


def get_ppt_from_cache(file_id: str) -> Optional[Dict[str, str]]:
    """从 Redis 缓存中获取 PPT 文件"""
    try:
        from app.database.redis_utils import redis_utils
        value = redis_utils.get(f"{PPT_CACHE_PREFIX}{file_id}")
        if value:
            return json.loads(value)
        return None
    except Exception as e:
        logger.error(f"PPT 从 Redis 读取失败: {e}", exc_info=True)
        return None


def _parse_color(color_str: str) -> Optional[RGBColor]:
    """解析颜色字符串，支持 #RRGGBB、RRGGBB(6位十六进制) 和常见颜色名"""
    if not color_str:
        return None
    color_str = color_str.strip()
    # 支持 #RRGGBB 格式
    if color_str.startswith("#") and len(color_str) == 7:
        try:
            return RGBColor(int(color_str[1:3], 16), int(color_str[3:5], 16), int(color_str[5:7], 16))
        except ValueError:
            return None
    # 支持 6 位十六进制字符串（无 # 前缀），如 "FF5733"
    if len(color_str) == 6 and all(c in '0123456789abcdefABCDEF' for c in color_str):
        try:
            return RGBColor(int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))
        except ValueError:
            return None
    color_map = {
        "black": RGBColor(0, 0, 0), "white": RGBColor(255, 255, 255),
        "red": RGBColor(255, 0, 0), "green": RGBColor(0, 128, 0),
        "blue": RGBColor(0, 0, 255), "yellow": RGBColor(255, 255, 0),
        "orange": RGBColor(255, 165, 0), "purple": RGBColor(128, 0, 128),
        "gray": RGBColor(128, 128, 128), "grey": RGBColor(128, 128, 128),
    }
    return color_map.get(color_str.lower())


def _set_text_style(run, style: Dict[str, Any]):
    """设置文本运行样式，支持 python-pptx Font 的所有常用属性"""
    if not isinstance(style, dict):
        return
    font = run.font
    if "font_size" in style:
        font.size = Pt(int(style["font_size"]))
    if "bold" in style:
        font.bold = bool(style["bold"])
    if "italic" in style:
        font.italic = bool(style["italic"])
    if "underline" in style:
        font.underline = bool(style["underline"])
    if "strikethrough" in style:
        font.strikethrough = bool(style["strikethrough"])
    if "color" in style:
        c = _parse_color(style["color"])
        if c:
            font.color.rgb = c
    if "font_name" in style:
        font.name = style["font_name"]
    if "font_name_east_asian" in style:
        try:
            font._element.rPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}ea', style["font_name_east_asian"])
        except Exception:
            pass
    if "shadow" in style:
        font.shadow = bool(style["shadow"])
    if "highlight_color" in style:
        c = _parse_color(style["highlight_color"])
        if c:
            font.highlight_color = c
    if "spacing" in style:
        font.spacing = Pt(float(style["spacing"]))
    if "superscript" in style and style["superscript"]:
        font.superscript = True
    if "subscript" in style and style["subscript"]:
        font.subscript = True


def _add_slide_from_definition(prs: Presentation, slide_def: Dict[str, Any], style_config: Dict[str, Any]):
    """根据幻灯片定义添加一页"""
    layout_index = slide_def.get("layout_index", 5)  # 默认空白布局
    try:
        layout = prs.slide_layouts[min(layout_index, len(prs.slide_layouts) - 1)]
    except IndexError:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # ---- 幻灯片标题 ----
    slide_title = slide_def.get("title", "")
    if slide_title:
        # 优先使用占位符标题
        title_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:  # Title
                title_shape = shape
                break
        if title_shape is None:
            # 手动添加标题文本框
            left = Inches(0.5)
            top = Inches(0.3)
            width = prs.slide_width - Inches(1)
            height = Inches(0.8)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = title_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.alignment = PP_ALIGN.LEFT
        title_style = style_config.get("title_style", {})
        for run in p.runs:
            _set_text_style(run, title_style)
        if not p.runs and title_style:
            run = p.add_run()
            run.text = slide_title
            p.clear()
            p.add_run().text = slide_title
            _set_text_style(p.runs[0], title_style)

    # ---- 副标题 ----
    subtitle = slide_def.get("subtitle", "")
    if subtitle:
        subtitle_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle / Body
                subtitle_shape = shape
                break
        if subtitle_shape is None:
            left = Inches(0.5)
            top = Inches(1.2)
            width = prs.slide_width - Inches(1)
            height = Inches(0.6)
            subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = subtitle_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.LEFT

    # ---- 正文内容区 ----
    content = slide_def.get("content", [])
    if content:
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Body
                body_shape = shape
                break
        if body_shape is None:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            height = prs.slide_height - Inches(2)
            body_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = body_shape.text_frame
        tf.word_wrap = True
        content_style = style_config.get("content_style", {})
        if not isinstance(content_style, dict):
            content_style = {}

        if isinstance(content, str):
            # 纯文本
            tf.paragraphs[0].text = content
            for run in tf.paragraphs[0].runs:
                _set_text_style(run, content_style)
        elif isinstance(content, list):
            first = True
            for item in content:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()

                if isinstance(item, str):
                    p.text = item
                    p.level = 0
                    for run in p.runs:
                        _set_text_style(run, content_style)
                elif isinstance(item, dict):
                    p.text = item.get("text", "")
                    p.level = item.get("level", 0)
                    item_style = item.get("style", {})
                    if not isinstance(item_style, dict):
                        item_style = {}
                    merged_style = {**content_style, **item_style}
                    for run in p.runs:
                        _set_text_style(run, merged_style)
                    if "alignment" in item:
                        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                        p.alignment = align_map.get(item["alignment"], PP_ALIGN.LEFT)

    # ---- 备注 ----
    notes = slide_def.get("notes", "")
    if notes:
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


@ToolRegistry.register
class generate_ppt(BaseTool):
    """基于 python-pptx 的 PPT 生成工具。"""

    name = "generate_ppt"
    title = "PPT生成"
    description = (
        "生成 PowerPoint (.pptx) 演示文稿。当用户需要创建幻灯片、演示文稿、PPT时使用此工具。"
        "支持自定义模板、多页幻灯片、标题/副标题/正文/列表等内容。"
        "生成成功后会返回下载链接(download_url)，请在回复中用 markdown 链接格式输出该下载链接供用户点击下载。"
    )
    params = [
        BaseToolParam(
            name="title",
            type="string",
            description="PPT 文档标题",
            required=True,
        ),
        BaseToolParam(
            name="slides",
            type="array",
            description=(
                "幻灯片列表，每项是一页幻灯片的定义对象，包含："
                "title(页面标题)、subtitle(副标题)、content(内容，字符串数组或对象数组，"
                "对象格式：{text, level, style})、notes(备注)、layout_index(布局索引，0-10，默认5)"
            ),
            required=True,
        ),
        BaseToolParam(
            name="template_ppt",
            type="string",
            description="PPT 模板文件，支持 base64 编码字符串或服务器上的文件路径。留空则使用默认空白模板。",
            required=False,
            default="",
        ),
        BaseToolParam(
            name="style_config",
            type="object",
            description=(
                "全局样式配置对象。包含 title_style 和 content_style 两个子对象，"
                "每个子对象的字段对应 python-pptx Font 属性，常用字段："
                "font_size(字号磅值)、bold(加粗)、italic(斜体)、underline(下划线)、"
                "strikethrough(删除线)、color(颜色,#RRGGBB格式)、font_name(字体名)、"
                "font_name_east_asian(东亚字体)、shadow(阴影)、highlight_color(高亮色)、"
                "spacing(字符间距磅值)。"
                "另支持 slide_width_inches 和 slide_height_inches 设置幻灯片尺寸。"
                "颜色尽量用 #RRGGBB 格式（如 #FF5733）"
            ),
            required=False,
        ),
        BaseToolParam(
            name="author",
            type="string",
            description="PPT 作者名称",
            required=False,
            default="",
        ),
        BaseToolParam(
            name="subject",
            type="string",
            description="PPT 主题/描述",
            required=False,
            default="",
        ),
        BaseToolParam(
            name="slide_width",
            type="number",
            description="幻灯片宽度（英寸），默认 13.33（宽屏16:9）",
            required=False,
        ),
        BaseToolParam(
            name="slide_height",
            type="number",
            description="幻灯片高度（英寸），默认 7.5",
            required=False,
        ),
    ]

    def run(self, **kwargs) -> Any:
        title = kwargs.get("title", "未命名演示文稿")
        slides = kwargs.get("slides", [])
        template_ppt = kwargs.get("template_ppt", "")
        style_config = kwargs.get("style_config", {}) or {}
        author = kwargs.get("author", "")
        subject = kwargs.get("subject", "")
        slide_width = kwargs.get("slide_width")
        slide_height = kwargs.get("slide_height")

        # 解析 slides（模型可能传 JSON 字符串）
        if isinstance(slides, str):
            try:
                slides = json.loads(slides)
            except json.JSONDecodeError:
                return "PPT 生成失败：slides 参数格式错误，需要 JSON 数组"
        if not isinstance(slides, list) or len(slides) == 0:
            return "PPT 生成失败：slides 不能为空，需要提供至少一页幻灯片定义"

        # 解析 style_config
        if isinstance(style_config, str):
            try:
                style_config = json.loads(style_config)
            except json.JSONDecodeError:
                style_config = {}

        # 覆盖宽高
        if slide_width:
            style_config["slide_width_inches"] = float(slide_width)
        if slide_height:
            style_config["slide_height_inches"] = float(slide_height)

        # ---- 创建 Presentation ----
        try:
            if template_ppt:
                prs = self._load_template(template_ppt)
            else:
                prs = Presentation()
        except Exception as e:
            logger.error(f"加载模板失败: {e}", exc_info=True)
            return f"PPT 生成失败：加载模板失败 - {str(e)}"

        # 设置幻灯片尺寸
        if "slide_width_inches" in style_config:
            prs.slide_width = Inches(float(style_config["slide_width_inches"]))
        if "slide_height_inches" in style_config:
            prs.slide_height = Inches(float(style_config["slide_height_inches"]))

        # ---- 添加幻灯片 ----
        for i, slide_def in enumerate(slides):
            if isinstance(slide_def, str):
                slide_def = {"title": slide_def, "content": []}
            try:
                _add_slide_from_definition(prs, slide_def, style_config)
            except Exception as e:
                logger.warning(f"第 {i+1} 页幻灯片生成失败: {e}", exc_info=True)

        # ---- 设置文档属性 ----
        prs.core_properties.title = title
        if author:
            prs.core_properties.author = author
        if subject:
            prs.core_properties.subject = subject

        # ---- 导出为 base64 ----
        try:
            file_id = uuid.uuid4().hex[:16]
            safe_title = title.replace("/", "_").replace("\\", "_")[:50]
            file_name = f"{safe_title}_{file_id}.pptx"

            # 写入内存缓冲区，生成 base64
            buf = BytesIO()
            prs.save(buf)
            buf.seek(0)
            file_base64 = base64.b64encode(buf.read()).decode("utf-8")

            # 存入 Redis 缓存，供下载端点使用
            save_ppt_to_cache(file_id, file_base64, file_name)

            download_url = f"/aicenter/v1/chat/download_ppt/{file_id}"

            return {
                "type": "ppt_file",
                "file_name": file_name,
                "file_id": file_id,
                "title": title,
                "slide_count": len(slides),
                "download_url": download_url,
                "message": f"成功生成 PPT「{title}」，共 {len(slides)} 页幻灯片。请使用 markdown 链接格式输出下载链接：[{file_name}]({download_url})"
            }
        except Exception as e:
            logger.error(f"导出 PPT 失败: {e}", exc_info=True)
            return f"PPT 生成失败：导出文件失败 - {str(e)}"

    @staticmethod
    def _load_template(template_ppt: str) -> Presentation:
        """加载模板：支持文件路径和 base64"""
        # 尝试作为文件路径
        if os.path.isfile(template_ppt):
            return Presentation(template_ppt)
        # 尝试作为 base64
        try:
            raw = base64.b64decode(template_ppt)
            return Presentation(BytesIO(raw))
        except Exception as e:
            raise ValueError(f"无法解析模板文件（既不是有效路径也不是有效 base64）: {str(e)}")
