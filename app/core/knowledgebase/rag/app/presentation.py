"""
演示文稿(PPT)文档切片方法
支持PPT、PPTX、PDF等演示文稿的按页切片
"""

import copy
import re
import logging
from io import BytesIO

try:
    from PIL import Image
except ImportError:
    Image = None

logger = logging.getLogger(__name__)


class Pdf:
    """PDF解析器（用于PPT转PDF后的处理）"""
    
    def __init__(self):
        pass
    
    def __call__(self, filename, binary=None, from_page=0, to_page=100000, zoomin=3, callback=None, **kwargs):
        """
        解析PDF格式的演示文稿
        
        每页作为一个chunk，并提取页面缩略图
        """
        try:
            from app.core.knowledgebase.deepdoc.parser import PdfParser
            
            pdf_parser = PdfParser()
            
            # 1. OCR识别
            callback(msg="开始OCR识别")
            pdf_parser.__images__(filename if not binary else binary, zoomin, from_page, to_page, callback)
            
            # 2. 布局分析
            callback(msg="布局分析")
            pdf_parser._layouts_rec(zoomin)
            
            # 3. 表格分析
            callback(msg="表格分析")
            pdf_parser._table_transformer_job(zoomin)
            
            # 4. 文本合并
            pdf_parser._text_merge()
            
            # 5. 提取表格和图片
            tbls = pdf_parser._extract_table_figure(True, zoomin, True, True)
            
            # 6. 按页重组内容
            from collections import defaultdict
            page_items = defaultdict(list)
            
            # 添加文本内容
            for b in pdf_parser.boxes:
                global_page_num = b["page_number"] + from_page
                if not (from_page < global_page_num <= to_page + from_page):
                    continue
                page_items[global_page_num].append({
                    "top": b["top"],
                    "x0": b["x0"],
                    "text": b["text"],
                    "type": "text"
                })
            
            # 添加表格和图片
            for (img, content), positions in tbls:
                if not positions:
                    continue
                    
                if isinstance(content, list):
                    final_text = "\n".join(content)
                elif isinstance(content, str):
                    final_text = content
                else:
                    final_text = str(content)
                
                try:
                    pn_index = positions[0][0]
                    if isinstance(pn_index, list):
                        pn_index = pn_index[0]
                    current_page_num = int(pn_index) + 1
                except Exception as e:
                    logger.error(f"解析位置错误: {e}")
                    continue
                
                if not (from_page < current_page_num <= to_page + from_page):
                    continue
                
                top = positions[0][3]
                left = positions[0][1]
                
                page_items[current_page_num].append({
                    "top": top,
                    "x0": left,
                    "text": final_text,
                    "type": "table_or_figure"
                })
            
            # 7. 生成结果
            res = []
            for i in range(len(pdf_parser.page_images)):
                current_pn = from_page + i + 1
                items = page_items.get(current_pn, [])
                items.sort(key=lambda x: (x["top"], x["x0"]))
                
                full_page_text = "\n\n".join([item["text"] for item in items])
                if not full_page_text.strip():
                    full_page_text = f"[第{current_pn}页无文本或数据]"
                    
                page_img = pdf_parser.page_images[i] if Image else None
                res.append((full_page_text, page_img))
            
            callback(0.9, "解析完成")
            return res, []
            
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            return [], []


def chunk(filename, binary=None, from_page=0, to_page=100000, lang="Chinese", callback=None, parser_config=None, **kwargs):
    """
    演示文稿切片函数
    
    支持PPT、PPTX、PDF格式，每页/每张幻灯片作为一个chunk
    
    Args:
        filename: 文件名或文件路径
        binary: 文件二进制数据（可选）
        from_page: 起始页码
        to_page: 结束页码
        lang: 语言 ("Chinese" 或 "English")
        callback: 进度回调函数 callback(progress, message)
        parser_config: 解析配置（可选）
        **kwargs: 其他参数
        
    Returns:
        list: 切片后的文档列表
    """
    from ..nlp import rag_tokenizer, tokenize
    
    if parser_config is None:
        parser_config = {}
        
    eng = lang.lower() == "english"
    doc = {"docnm_kwd": filename}
    doc["title_tks"] = rag_tokenizer.tokenize(re.sub(r"\.[a-zA-Z]+$", "", filename))
    doc["title_sm_tks"] = rag_tokenizer.fine_grained_tokenize(doc["title_tks"])
    res = []
    
    # PPT/PPTX文件处理
    if re.search(r"\.pptx?$", filename, re.IGNORECASE):
        callback(0.1, "开始解析PPT/PPTX")
        
        try:
            sections = _parse_pptx(filename, binary, from_page, to_page, callback)
            
            for pn, (txt, img) in enumerate(sections):
                d = copy.deepcopy(doc)
                pn += from_page
                d["doc_type_kwd"] = "image"
                d["page_num_int"] = [pn + 1]
                d["top_int"] = [0]
                d["position_int"] = [(pn + 1, 0, 0, 0, 0)]
                
                if img and Image and isinstance(img, Image.Image):
                    d["image"] = img
                
                tokenize(d, txt, eng)
                res.append(d)
            
            callback(0.8, f"PPT解析完成，共 {len(res)} 页")
            return res
            
        except Exception as e:
            logger.warning(f"python-pptx解析失败: {e}，尝试其他方式")
            if callback:
                callback(0.2, "python-pptx失败，尝试其他方式")
    
    # PDF文件处理
    elif re.search(r"\.pdf$", filename, re.IGNORECASE):
        callback(0.1, "开始解析PDF演示文稿")
        
        layout_recognizer = parser_config.get("layout_recognize", "DeepDOC")
        if isinstance(layout_recognizer, bool):
            layout_recognizer = "DeepDOC" if layout_recognizer else "Plain Text"
        
        pdf_parser = Pdf()
        sections, _ = pdf_parser(
            filename=filename,
            binary=binary,
            from_page=from_page,
            to_page=to_page,
            callback=callback,
            **kwargs
        )
        
        if not sections:
            return []
        
        for pn, (txt, img) in enumerate(sections):
            d = copy.deepcopy(doc)
            pn += from_page
            
            if img and Image and isinstance(img, Image.Image):
                d["image"] = img
            else:
                d["image"] = None
                
            d["page_num_int"] = [pn + 1]
            d["top_int"] = [0]
            d["position_int"] = [
                (pn + 1, 0, img.size[0] if img and hasattr(img, 'size') else 0, 
                 0, img.size[1] if img and hasattr(img, 'size') else 0)
            ]
            
            tokenize(d, txt, eng)
            res.append(d)
        
        callback(0.8, f"PDF演示文稿解析完成，共 {len(res)} 页")
        return res
    
    raise NotImplementedError(f"不支持的文件类型: {filename} (支持: ppt, pptx, pdf)")


def _parse_pptx(filename, binary=None, from_page=0, to_page=1000000, callback=None):
    """解析PPTX文件"""
    try:
        from pptx import Presentation
        from io import BytesIO
    except ImportError:
        raise ImportError("请安装python-pptx: pip install python-pptx")
    
    if binary:
        prs = Presentation(BytesIO(binary))
    else:
        prs = Presentation(filename)
    
    sections = []
    
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx < from_page:
            continue
        if slide_idx >= to_page:
            break
            
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        
        full_text = "\n".join(texts)
        
        # 尝试提取缩略图
        thumbnail = None
        if Image:
            try:
                # 使用PIL创建一个空白图像作为占位符
                # 实际应用中可以使用pptx的图片提取功能
                pass
            except:
                pass
        
        sections.append((full_text, thumbnail))
        
        if callback and (slide_idx + 1) % 10 == 0:
            callback(slide_idx / len(prs.slides), f"正在解析第 {slide_idx + 1}/{len(prs.slides)} 页")
    
    return sections


if __name__ == "__main__":
    import sys
    
    def dummy_callback(prog=None, msg=""):
        print(f"[{prog}] {msg}" if prog else msg)
    
    if len(sys.argv) > 1:
        result = chunk(sys.argv[1], callback=dummy_callback)
        print(f"Total chunks: {len(result)}")
